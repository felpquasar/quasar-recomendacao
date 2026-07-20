# -*- coding: utf-8 -*-
"""
Gera slides.pptx da apresentacao de defesa do TCC.
Sistema de Recomendacao B2B com Regras de Associacao - Quasar Barber.
Numeros verificados contra dados/metricas-validacao.json e algoritmo real.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Tema ----
NAVY   = RGBColor(0x0F, 0x1B, 0x2D)   # fundo escuro / titulos
INK     = RGBColor(0x1A, 0x23, 0x2E)   # texto corpo
GOLD    = RGBColor(0xC8, 0xA0, 0x4B)   # acento
LIGHT   = RGBColor(0xF4, 0xF1, 0xEA)   # fundo claro
GRAY    = RGBColor(0x5A, 0x66, 0x72)   # texto secundario
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(sh, color)
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h, lines, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, space=6):
    """lines: str ou lista de (texto, dict_opcoes)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        opts = {}
        if isinstance(ln, tuple):
            ln, opts = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space", space))
        p.space_before = Pt(0)
        # suporte a bullet com runs multiplos via lista
        runs = ln if isinstance(ln, list) else [(ln, {})]
        for txt, ro in runs:
            r = p.add_run()
            r.text = txt
            r.font.name = ro.get("font", opts.get("font", font))
            r.font.size = Pt(ro.get("size", opts.get("size", size)))
            r.font.bold = ro.get("bold", opts.get("bold", bold))
            r.font.color.rgb = ro.get("color", opts.get("color", color))
    return tb


def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, LIGHT)
    # barra lateral acento
    rect(s, 0, 0, Inches(0.18), SH, GOLD)
    # titulo
    textbox(s, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.9),
            title, size=30, color=NAVY, bold=True)
    if kicker:
        textbox(s, Inches(0.62), Inches(0.18), Inches(12), Inches(0.35),
                kicker.upper(), size=12, color=GOLD, bold=True)
    # linha
    rect(s, Inches(0.62), Inches(1.32), Inches(2.4), Pt(3), GOLD)
    # rodape: autor (esq) + numero do slide (dir)
    num = len(prs.slides._sldIdLst)
    textbox(s, Inches(0.62), Inches(7.05), Inches(9), Inches(0.35),
            "Felipe Pereira  •  Sistema de Recomendação B2B com Regras de Associação",
            size=10, color=GRAY)
    textbox(s, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.35),
            str(num), size=11, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
    return s


def bullets(slide, items, x=Inches(0.7), y=Inches(1.7), w=Inches(12), h=Inches(5.3),
            size=20, gap=14):
    lines = []
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        prefix = "•  " if lvl == 0 else "–  "
        col = INK if lvl == 0 else GRAY
        sz = size if lvl == 0 else size - 3
        bold_seg = []
        # negrito em **...**
        import re
        parts = re.split(r"(\*\*.*?\*\*)", it)
        runs = [(prefix, {"color": GOLD, "bold": True, "size": sz})]
        for pt in parts:
            if pt.startswith("**") and pt.endswith("**"):
                runs.append((pt[2:-2], {"color": NAVY, "bold": True, "size": sz}))
            elif pt:
                runs.append((pt, {"color": col, "size": sz}))
        indent = Inches(0.4) if lvl else Inches(0)
        lines.append((runs, {"space": gap}))
    textbox(slide, x, y, w, h, lines, space=gap)


# =================================================================
# SLIDE 1 - CAPA
# =================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(5.55), SW, Inches(0.12), GOLD)
textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.8),
        [("[UNIVERSIDADE / INSTITUIÇÃO]", {"size": 15, "color": WHITE, "bold": True, "space": 2}),
         ("Curso de Ciência da Computação", {"size": 13, "color": RGBColor(0xC9,0xCF,0xD6)})])
textbox(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6),
        ["Sistema de Recomendação de Produtos para",
         "Distribuição B2B Utilizando Regras de Associação",
         ("Estudo de Caso na Quasar Barber", {"size": 22, "color": GOLD, "space": 0})],
        size=38, color=WHITE, bold=True, space=4)
textbox(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.3),
        [("Felipe Pereira", {"size": 20, "color": WHITE, "bold": True, "space": 2}),
         ("Trabalho de Conclusão de Curso  —  Ciência da Computação", {"size": 15, "color": RGBColor(0xC9,0xCF,0xD6)}),
         ("Orientador(a): [Nome]   |   Codó-MA, 2026", {"size": 13, "color": GRAY})])

# =================================================================
# SLIDE 2 - AGENDA
# =================================================================
s = content_slide("Roteiro", "Apresentação")
bullets(s, [
    "Contexto e problema de pesquisa",
    "Pergunta, hipótese e objetivos",
    "Fundamentação: regras de associação",
    "Metodologia e dados reais",
    "Arquitetura do sistema",
    "Protocolo de validação (LOOCV)",
    "Resultados e discussão",
    "Conclusão e trabalhos futuros",
], y=Inches(1.8), size=22, gap=12)

# =================================================================
# SLIDE 3 - CONTEXTO / PROBLEMA
# =================================================================
s = content_slide("Contexto e problema", "Capítulo 1")
bullets(s, [
    "Distribuidora B2B de **pequeno porte**: catálogo amplo, poucos clientes.",
    "Decisões de **cross-sell** (venda casada) tomadas por **intuição** do vendedor.",
    "Cenário B2B difere do varejo B2C:",
    ("Poucos clientes e baixo volume transacional", 1),
    ("Sem avaliações ou notas de produto", 1),
    ("Matriz usuário-item esparsa → filtragem colaborativa **inviável**", 1),
    "Oportunidade: minerar o **histórico de vendas** para sugerir produtos complementares.",
], y=Inches(1.8), size=21, gap=12)

# =================================================================
# SLIDE 4 - PERGUNTA + HIPOTESE
# =================================================================
s = content_slide("Pergunta de pesquisa e hipótese", "Capítulo 1")
# caixa pergunta
box = rect(s, Inches(0.7), Inches(1.75), Inches(12), Inches(1.7), NAVY)
textbox(s, Inches(1.0), Inches(1.9), Inches(11.4), Inches(1.4),
        ["É possível gerar recomendações úteis e estatisticamente relevantes",
         "para uma distribuidora B2B de pequeno porte usando apenas regras de",
         "associação, mesmo com volume reduzido de transações?"],
        size=19, color=WHITE, bold=True, space=2, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, [
    "**H1:** regras com confiança ≥ 0,30 e suporte ≥ 0,02, ranqueadas por lift, geram cross-sell coerente com o domínio.",
    "**H0:** as regras são indistinguíveis do acaso e não superam um baseline de popularidade.",
], y=Inches(3.8), size=19, gap=14)

# =================================================================
# SLIDE 5 - OBJETIVOS
# =================================================================
s = content_slide("Objetivos", "Capítulo 1")
textbox(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.6),
        [[("Geral:  ", {"color": GOLD, "bold": True, "size": 20}),
          ("desenvolver e validar um sistema de recomendação B2B baseado em Market Basket Analysis, com relevância estatística mensurável.", {"color": INK, "size": 20})]])
textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(0.4),
        "Específicos:", size=18, color=GOLD, bold=True)
bullets(s, [
    "Revisar a literatura (recomendação, Apriori, support/confidence/lift).",
    "Modelar base transacional (PostgreSQL/Supabase).",
    "Implementar extração de regras e recomendação Top-N (API REST).",
    "Construir interface web de consulta (React).",
    "Validar via LOOCV (Precision, Recall, F1).",
    "Analisar as regras de maior lift sob a ótica de negócio.",
], y=Inches(3.4), size=18, gap=9)

# =================================================================
# SLIDE - JUSTIFICATIVA
# =================================================================
s = content_slide("Justificativa", "Capítulo 1")
just = [
    ("Acadêmica", "Maioria dos estudos foca B2C de grande volume. Lacuna em B2B esparso, onde a filtragem colaborativa falha."),
    ("Prática", "Empresa real do autor. Apoia a equipe de vendas com sugestões baseadas em dados, não em intuição."),
    ("Técnica", "Arquitetura completa, reproduzível e de baixo custo, viável para pequenas empresas sem Big Data."),
]
y = Inches(1.9); bh = Inches(1.45); gap = Inches(0.3)
for i, (t, d) in enumerate(just):
    yy = y + i * (bh + gap)
    rect(s, Inches(0.7), yy, Inches(12), bh, WHITE)
    rect(s, Inches(0.7), yy, Inches(0.12), bh, GOLD)
    textbox(s, Inches(1.05), yy, Inches(2.7), bh, t, size=21, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(3.9), yy, Inches(8.5), bh, d, size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# =================================================================
# SLIDE 6 - FUNDAMENTACAO: METRICAS
# =================================================================
s = content_slide("Regras de associação", "Capítulo 2  •  Market Basket Analysis")
textbox(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.5),
        "Cada pedido é uma cesta; padrão A ⇒ B. Três métricas:", size=20, color=INK)
cards = [
    ("Support", "freq(A∩B) / total", "Quão comum é o par. Filtra raros."),
    ("Confidence", "freq(A∩B) / freq(A)", "Dado A, probabilidade de B. Força da regra."),
    ("Lift", "confidence / support(B)", "> 1 = associação real, não acaso."),
]
cw = Inches(3.95); gap = Inches(0.18); x0 = Inches(0.7); y0 = Inches(2.5); ch = Inches(2.7)
for i, (t, f, d) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE)
    rect(s, x, y0, cw, Inches(0.12), GOLD)
    textbox(s, x, y0 + Inches(0.35), cw, Inches(0.6), t, size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.2), y0 + Inches(1.1), cw - Inches(0.4), Inches(0.6), f,
            size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.25), y0 + Inches(1.75), cw - Inches(0.5), Inches(0.9), d,
            size=15, color=GRAY, align=PP_ALIGN.CENTER)

# =================================================================
# SLIDE 7 - POR QUE ASSOCIATION RULES
# =================================================================
s = content_slide("Por que regras de associação?", "Capítulo 2")
bullets(s, [
    "**Filtragem colaborativa** exige matriz usuário-item densa → falha em B2B esparso.",
    "**Cold start** e esparsidade são a regra, não a exceção, nesse cenário.",
    "Regras de associação operam direto sobre **coocorrência em transações**:",
    ("Não precisam de notas nem de muitos clientes", 1),
    ("Resultado interpretável pela equipe de vendas", 1),
    ("Baixo custo computacional no volume estudado", 1),
    "Lacuna acadêmica: maioria dos estudos foca B2C de **grande volume**.",
], y=Inches(1.8), size=20, gap=11)

# =================================================================
# SLIDE 8 - METODOLOGIA + DADOS
# =================================================================
s = content_slide("Metodologia e dados reais", "Capítulo 3")
bullets(s, [
    "Pesquisa **aplicada**, quantitativa, estudo de caso único.",
    "Método: **Design Science Research** (artefato + avaliação).",
    "Dados reais de venda da Quasar Barber: 29/04 a 18/06/2026.",
], y=Inches(1.7), w=Inches(12), size=20, gap=10)
# faixa de numeros
stats = [("30", "transações"), ("26", "clientes"), ("29", "produtos"),
         ("52", "itens"), ("14", "multi-item")]
sw = Inches(2.35); gap = Inches(0.12); x0 = Inches(0.7); y0 = Inches(4.4); ch = Inches(1.7)
for i, (n, lab) in enumerate(stats):
    x = x0 + i * (sw + gap)
    rect(s, x, y0, sw, ch, NAVY)
    textbox(s, x, y0 + Inches(0.25), sw, Inches(0.9), n, size=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x, y0 + Inches(1.15), sw, Inches(0.4), lab, size=15, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.5),
        "Apenas 14 cestas multi-item geram regras — a esparsidade é intrínseca ao B2B.",
        size=15, color=GRAY, align=PP_ALIGN.CENTER)

# =================================================================
# SLIDE 9 - ARQUITETURA
# =================================================================
s = content_slide("Arquitetura do sistema", "Capítulo 4")
# print real: diagrama de arquitetura em camadas (tcc/figuras)
arq_h = Inches(5.4); arq_w = Inches(5.4 * 1.357)
arq_x = (SW - arq_w) // 2
s.shapes.add_picture("tcc/figuras/figura1-arquitetura.png", arq_x, Inches(1.65), height=arq_h)

# =================================================================
# SLIDE 10 - VALIDACAO
# =================================================================
s = content_slide("Protocolo de validação", "Capítulo 3  •  por que LOOCV")
bullets(s, [
    "**Holdout 80/20 descartado:** em 30 transações, sobram só **3** cestas avaliáveis no teste → instável.",
    "**LOOCV (Leave-One-Out):** N rodadas, treina com N-1, testa 1. Aproveita todos os dados.",
    "Esquema **esconder-1-item**: oculta um produto da cesta, sistema tenta recuperá-lo no Top-3.",
    "**Teto da Precision = 1/N ≈ 0,33:** 1 produto a achar, até 3 recomendados.",
    ("→ Recall e F1 são os indicadores mais informativos.", 1),
], y=Inches(1.8), size=19, gap=14)

# =================================================================
# SLIDE 11 - RESULTADOS: REGRAS
# =================================================================
s = content_slide("Resultados: regras descobertas", "Capítulo 5")
LCOL = Inches(8.4)   # largura coluna esquerda (texto + tabela)
textbox(s, Inches(0.7), Inches(1.55), LCOL, Inches(0.9),
        [[("44 regras geradas", {"color": GOLD, "bold": True, "size": 21})],
         [("confiança média 76,51%  •  lift médio 12,41", {"color": INK, "size": 17})]], space=2)
# tabela top regras (coluna esquerda)
rows = [
    ("Antecedente ⇒ Consequente", "Conf.", "Lift"),
    ("After Shave (QOD) ⇒ Café Verde Shampoo", "1,00", "30,0"),
    ("Silver Boost (QOD) ⇒ Bálsamo Fortalecedor", "1,00", "30,0"),
    ("Shampoo ⇒ Óleo Barba (Fox for Men)", "1,00", "15,0"),
    ("Maquiagem ⇒ Aplicador Fibra (Fox)", "1,00", "15,0"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(2.6), LCOL, Inches(2.8)).table
tbl.columns[0].width = Inches(5.8)
tbl.columns[1].width = Inches(1.3); tbl.columns[2].width = Inches(1.3)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(14 if ri else 13)
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE if ri == 0 else INK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (LIGHT if ri % 2 else WHITE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
textbox(s, Inches(0.7), Inches(5.6), LCOL, Inches(1.0),
        "Lift alto coincide com baixo suporte (pares raros): interpretar junto ao domínio.",
        size=14, color=GRAY)
# print real do dashboard (coluna direita)
dash_h = Inches(5.2); dash_w = Inches(5.2 * 0.711)
dash_x = Inches(9.3)
s.shapes.add_picture("tcc/figuras/figura3-dashboard.png", dash_x, Inches(1.65), height=dash_h)
textbox(s, dash_x, Inches(6.95), dash_w, Inches(0.35),
        "Painel do sistema", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# =================================================================
# SLIDE 12 - RESULTADOS: VALIDACAO
# =================================================================
s = content_slide("Resultados: validação LOOCV", "Capítulo 5")
metrics = [("Precision", "26,39%", "próx. ao teto 0,33"),
           ("Recall", "44,44%", "16 de 36 sub-testes"),
           ("F1-Score", "31,94%", "média harmônica")]
cw = Inches(3.95); gap = Inches(0.18); x0 = Inches(0.7); y0 = Inches(1.9); ch = Inches(2.3)
for i, (t, v, d) in enumerate(metrics):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, NAVY)
    textbox(s, x, y0 + Inches(0.25), cw, Inches(0.5), t, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x, y0 + Inches(0.85), cw, Inches(0.8), v, size=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x, y0 + Inches(1.7), cw, Inches(0.5), d, size=14, color=RGBColor(0xC9,0xCF,0xD6), align=PP_ALIGN.CENTER)
bullets(s, [
    "14 transações avaliadas, 36 sub-testes acumulados.",
    "Recall 44,44%: cesta completada em quase metade dos casos, sem notas de produto.",
    "Holdout 80/20 daria só 22,22% sobre 3 casos → LOOCV é mais robusto.",
], y=Inches(4.6), size=18, gap=10)

# =================================================================
# SLIDE 13 - ADERENCIA + HIPOTESE
# =================================================================
s = content_slide("Aderência ao domínio e hipótese", "Capítulo 5")
bullets(s, [
    "Regras de maior lift fazem **sentido comercial**:",
    ("Rotina de barba: shampoo + óleo + bálsamo + pós-barba", 1),
    ("Coerência de marca: linhas QOD e Fox for Men", 1),
    ("Acessórios de função equivalente compr. juntos", 1),
], y=Inches(1.7), size=19, gap=10)
box = rect(s, Inches(0.7), Inches(4.4), Inches(12), Inches(2.1), NAVY)
textbox(s, Inches(1.0), Inches(4.6), Inches(11.4), Inches(1.8),
        [[("H1 não rejeitada.  ", {"color": GOLD, "bold": True, "size": 22}),
          ("Lift médio 12,41 ≫ 1 afasta o acaso; pares de maior lift", {"color": WHITE, "size": 18})],
         ("não coincidem com os mais vendidos → carregam informação além da", {"color": WHITE, "size": 18, "space": 2}),
         ("frequência, contrariando H0 (baseline de popularidade).", {"color": WHITE, "size": 18, "space": 2})],
        anchor=MSO_ANCHOR.MIDDLE, space=2)

# =================================================================
# SLIDE 14 - CONCLUSAO
# =================================================================
s = content_slide("Conclusão", "Capítulo 6")
bullets(s, [
    "Sim: regras de associação geram recomendações relevantes em B2B de baixo volume — **com ressalva de escala**.",
    "Artefato **completo e reproduzível**, de baixo custo, sobre dados reais.",
    "Contribuição metodológica: **LOOCV** + leitura honesta do **teto da Precision**.",
    "Recomendações como **apoio à decisão**, não prescrição automática.",
], y=Inches(1.8), size=20, gap=14)

# =================================================================
# SLIDE 15 - TRABALHOS FUTUROS
# =================================================================
s = content_slide("Trabalhos futuros", "Capítulo 6")
bullets(s, [
    "**Mais dados:** histórico maior → métricas mais estáveis.",
    "**Apriori / FP-Growth completos:** regras multi-item (3+ produtos).",
    "**Deploy em produção:** recomendação no momento do pedido.",
    "**Teste A/B:** medir efeito real no ticket médio.",
], y=Inches(1.9), size=22, gap=18)

# =================================================================
# SLIDE - REFERENCIAS PRINCIPAIS
# =================================================================
s = content_slide("Referências principais", "ABNT NBR 6023")
refs = [
    "AGRAWAL, R.; SRIKANT, R. Fast algorithms for mining association rules. In: VLDB, 1994. p. 487-499.",
    "AGRAWAL, R.; IMIELIŃSKI, T.; SWAMI, A. Mining association rules between sets of items in large databases. In: ACM SIGMOD, 1993. p. 207-216.",
    "HAN, J.; PEI, J.; YIN, Y. Mining frequent patterns without candidate generation. In: ACM SIGMOD, 2000. p. 1-12.",
    "ADOMAVICIUS, G.; TUZHILIN, A. Toward the next generation of recommender systems. IEEE TKDE, v. 17, n. 6, p. 734-749, 2005.",
    "SARWAR, B. et al. Item-based collaborative filtering recommendation algorithms. In: WWW, 2001. p. 285-295.",
    "HAN, J.; KAMBER, M.; PEI, J. Data mining: concepts and techniques. 3. ed. Morgan Kaufmann, 2011.",
    "HEVNER, A. R. et al. Design science in information systems research. MIS Quarterly, v. 28, n. 1, p. 75-105, 2004.",
    "POWERS, D. M. W. Evaluation: from precision, recall and F-measure to ROC... J. Machine Learning Tech., v. 2, n. 1, 2011.",
]
lines = []
for r in refs:
    lines.append(([("•  ", {"color": GOLD, "bold": True, "size": 13}),
                   (r, {"color": INK, "size": 13})], {"space": 10}))
textbox(s, Inches(0.7), Inches(1.75), Inches(12), Inches(5.1), lines, space=10)
textbox(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
        "Lista completa (23 referências) no capítulo Referências da monografia.",
        size=12, color=GRAY)

# =================================================================
# SLIDE 16 - ENCERRAMENTO
# =================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0), Inches(3.5), SW, Inches(0.1), GOLD)
textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.0),
        "Obrigado.", size=48, color=WHITE, bold=True)
textbox(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(1.2),
        [("Felipe Pereira", {"size": 22, "color": GOLD, "bold": True, "space": 2}),
         ("Sistema de Recomendação B2B com Regras de Associação", {"size": 16, "color": RGBColor(0xC9,0xCF,0xD6)}),
         ("Quasar Barber  •  Codó-MA  •  2026", {"size": 14, "color": GRAY})])

prs.save("apresentacao/slides.pptx")
print("OK -", len(prs.slides._sldIdLst), "slides ->", "apresentacao/slides.pptx")
